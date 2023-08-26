r'''
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from imagesize import imagesize
from removetrashes import removeTrash
from namegenerator import nameGenerator
import glob
import time
import openai
import re

root = Presentation()
SLIDE_LAYOUT = root.slide_layouts[6] 

def addnewlines(text,LINE_LEN=90):
    text = text.replace('\n','')
    cutcount = 0
    for i in range(1,len(text)//LINE_LEN+1):
        text = text[:i*LINE_LEN+i-1]+'\n'+ text[i*LINE_LEN+i-1:]
        cutcount+=1
    return text,cutcount

def first_page(title,name,teacher,grade):
    global root
    slide = root.slides.add_slide(SLIDE_LAYOUT)

    # shir logo
    imgpath = r'C:\Users\User\Downloads\shirlogo.jpg'
    img = slide.shapes.add_picture(
                imgpath,
                Inches(0.01),Inches(0.35))

    # title 
    top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(Inches(0.35),top, width, height)  
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title.capitalize()
    p.font.bold = True
    p.font.underline = True
    p.font.textshadow = True # how to add shadow?
    p.font.size = Pt(45) #add addnewlines function for him

    # information
    top = width = height = Inches(3)
    txBox = slide.shapes.add_textbox(Inches(0.35), top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f'{time.asctime()}\n'
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = f'name {name}'
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = f'teacher {teacher}'
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = f'grade {grade}'
    p.font.size = Pt(24)
    
openai.api_key = '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'

def getcontent(message):
    msgs = []
    msgs.append(
            {'role':'user',
             'content':message})

    response = openai.ChatCompletion.create(
        model='gpt-3.5-turbo-16k',
        messages=msgs,
        temperature=0.2,
        max_tokens=15000,
    )
    content = response['choices'][0]['message']['content']
    return content,'text'

def last_page(links):
    global root        
    slide = root.slides.add_slide(SLIDE_LAYOUT)
    
    # title
    left = top = width = height = Inches(0.35)
    txBox = slide.shapes.add_textbox(left,top, width, height)  
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Links\n' # language correction
    p.font.bold = True
    p.font.underline = True
    p.font.textshadow = True # how to add shadow?
    p.font.size = Pt(45) #add addnewlines function for him

    # links
    for index,line in enumerate(links):
        p = tf.add_paragraph()
        if line[0]=='"':
            p.text = ' >' + addnewlines(line,90)[0]
        else:
            p.text = addnewlines(line,90)[0]
        p.font.size = Pt(16)
        if index!=0:
            p.font.color.rgb = RGBColor(26,115,232)
            p.font.underline = True
    
    savename = nameGenerator() + '.pptx'
    root.save(savename)
    return savename

def slidegptmaker(data,language,globsearch=r'C:\Users\User\Documents\html\задачник по html и css\pictures\*'):
    global root 
    imgurls = glob.glob(globsearch)
    imgcount = 0
    if len(data[0])>len(data[1]) or '[' in data[3]:
        data = data[1:]
    
    while len(data)>=3:
        data = removeTrash(data,'')
        title = data[0]
        title = translateTo(title,language)
        del data[0]

        description = data[0]
        description = translateTo(description,language)
        del data[0]

        picture = data[0]
        del data[0]

        # Create Title
        slide = root.slides.add_slide(SLIDE_LAYOUT)
        left = top = width = height = Inches(0.3)
        txBox = slide.shapes.add_textbox(left,top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.underline = True
        p.font.textshadow = True # how to add shadow?
        p.font.size = Pt(40)

        #Create Description
        left = top = width = height = Inches(1.25)
        txBox = slide.shapes.add_textbox(Inches(0.35), top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text,cutcount = addnewlines(description)
        p.font.size = Pt(18)
        
        # Create Picture
        try:
            imgurl = imgurls[imgcount]
            imgpath = imagesize(
                imgurl,
                cutcount)
            img = slide.shapes.add_picture(
                imgpath,
                Inches(0.35),Inches(1.25+(0.25*cutcount)+0.416+0.425))
        except:
            left = top = width = height = Inches(3)
            txBox = slide.shapes.add_textbox(Inches(0.35), top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text,cutcount = addnewlines(picture)
            p.font.size = Pt(18)
        imgcount+=1
    savename = nameGenerator() + '.pptx'
    root.save(savename)
    return savename

#print(last_page(['global warming','Aren','Arega','1036']))

arr = ['a','b','c']

def simbols(a,b,c='d'):
    print(a,b,c)

from gpt import getcontent
from wordcreater import wordCreater

wordCreater(getcontent(
Given a natural number n. 
Write a function that determines whether n is the square of a natural number, true/false.
in c++ use using namespace std;.
Dont use explanations.)[0].split('\n'),'en')

name = 'essaywriter'
dict = {
    'slidecreator': {
        'lang': True,
        'words': False,
        'title': True
    },
    'essaywriter': {
        'lang': True,
        'words': True,
        'title': True
    }
}

for item in dict:
    if item['title']:
        print('add title')

for item in dict:
    if item['words']:
        print('words')
'''

name = 'essaywriter'
names = ['essaywriter','slidecretor']

if name in names:
    print('add params')
