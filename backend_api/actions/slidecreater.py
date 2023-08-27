from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from backend_api.actions.imagesize import imagesize
from backend_api.actions.removetrashes import removeTrash
from backend_api.actions.namegenerator import nameGenerator
from backend_api.actions.translator import translateTo
import glob
import time
import os

root = Presentation()
SLIDE_LAYOUT = root.slide_layouts[6] 

def first_page(title,name,teacher,grade,language):
    global root
    slide = root.slides.add_slide(SLIDE_LAYOUT)

    # shir logo
    r'''
    imgpath = r'C:\Users\User\Downloads\shirlogo.jpg'
    img = slide.shapes.add_picture(
                imgpath,
                Inches(0.01),Inches(0.35))
    '''

    # title 
    top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(Inches(0.35),top, width, height)  
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text,tcut = addnewlines(translateTo(title,language)[0].capitalize(),28)
    p.font.bold = True
    p.font.underline = True
    # p.font.textshadow = True # how to add shadow?
    p.font.size = Pt(38) #add addnewlines function for him

    # information
    top = width = height = Inches(3)
    if tcut:
        top = Inches(4)
    txBox = slide.shapes.add_textbox(Inches(0.35), top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f'{time.asctime()}\n'
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = translateTo(f'name {name}',language)[0]
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = translateTo(f'teacher {teacher}',language)[0]
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = translateTo(f'grade {grade}',language)[0]
    p.font.size = Pt(24)

def last_page(links,language):
    global root        
    slide = root.slides.add_slide(SLIDE_LAYOUT)

    # title
    left = top = width = height = Inches(0.35)
    txBox = slide.shapes.add_textbox(left,top, width, height)  
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = translateTo('Links',language)[0] + '\n' # language correction
    p.font.bold = True
    p.font.underline = True
    # p.font.textshadow = True # how to add shadow?
    p.font.size = Pt(45) #add addnewlines function for him

    # links
    for index,line in enumerate(links):
        if index>=20:
            break
        p = tf.add_paragraph()
        if line[0]=='"':
            p.text = ' >' + addnewlines(line,70)[0][:150]
        else:
            p.text = addnewlines(line,70)[0][:150]
        p.font.size = Pt(16)
        if index!=0:
            p.font.color.rgb = RGBColor(26,115,232)
            p.font.underline = True
        
def slidegptmaker(data,links,language, 
                  globsearch=r'',
                  title = 'global warming',name = 'Name',
                  teacher = 'Teacher', grade = 'Grade',
                  ):
    global root 
    if globsearch:
        imgurls = glob.glob(globsearch)
    imgcount = 0
    if len(data[0])>len(data[1]) or '[' in data[3]:
        data = data[1:]
    
    first_page(title,name,teacher,grade,language)
    while len(data)>=3:
        data = removeTrash(data,'')
        title = data[0]
        title = translateTo(title,language)[0]
        del data[0]

        description = data[0]
        description = translateTo(description,language)[0]
        del data[0]

        picture = data[0]
        del data[0]

        # Create Title
        slide = root.slides.add_slide(SLIDE_LAYOUT)
        left = top = width = height = Inches(0.3)
        txBox = slide.shapes.add_textbox(left,top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text,tcut = addnewlines(title,28)
        p.font.bold = True
        p.font.underline = True
        p.font.textshadow = True # how to add shadow?
        p.font.size = Pt(40)

        #Create Description
        left = top = width = height = Inches(1.25)
        if tcut:
            top = Inches(1.25+(0.5*tcut))
        txBox = slide.shapes.add_textbox(Inches(0.35), top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text,cutcount = addnewlines(description,70)
        p.font.size = Pt(18)
        
        # Create Picture
        try:
            imgurl = imgurls[imgcount]
            imgpath = imagesize(
                imgurl,
                cutcount,
                tcut)
            img = slide.shapes.add_picture(
                imgpath,
                Inches(0.35),Inches(1.25+(0.25*cutcount)+(0.5*tcut)+0.416+0.425))
        except:
            left = top = width = height = Inches(3)
            txBox = slide.shapes.add_textbox(Inches(0.35), top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text,cutcount = addnewlines(picture,90)
            p.font.size = Pt(18)
        imgcount+=1
    last_page(links,language)
    
    name = nameGenerator() + '.pptx'
    savepath = os.path.join('media','pptx',name)
    root.save(savepath)
    return savepath
    
def addnewlines(text,line_len):
    text = text.replace('\n','')
    cutcount = 0
    for i in range(1,len(text)//line_len+1):
        text = text[:i*line_len+i-1]+'\n'+ text[i*line_len+i-1:]
        cutcount+=1
    return text,cutcount
